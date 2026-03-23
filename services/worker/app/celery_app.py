"""Worker - Celery tasks for document processing & URL fetching"""

from celery import Celery
import os
import sys
import hashlib
import httpx
from typing import List, Dict, Any
from io import BytesIO

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from common.config import Settings
from common.database import Database
from common.models import Document, IngestionTask, Chunk, FAQEntry, Entity
from common.document_status import (
    mark_document_failed,
    mark_parse_completed,
    mark_parse_processing,
    mark_summary_completed,
    mark_summary_processing,
)

settings = Settings()
db = Database(settings.postgres_url)

redis_url = os.getenv("REDIS_URL", "redis://redis:6379/0")
celery_app = Celery("worker", broker=redis_url, backend=redis_url)

celery_app.conf.update(
    task_serializer='json',
    accept_content=['json'],
    result_serializer='json',
    timezone='UTC',
    enable_utc=True,
)

from minio import Minio
minio_client = Minio(
    settings.minio_endpoint,
    access_key=settings.minio_access_key,
    secret_key=settings.minio_secret_key,
    secure=settings.minio_secure
)

from qdrant_client import QdrantClient
from qdrant_client.models import VectorParams, Distance, PointStruct
qdrant_client = QdrantClient(url=settings.qdrant_url)

COLLECTION_NAME = "kb_chunks_v2"


def ensure_collection_exists(vector_size: int):
    try:
        collections = qdrant_client.get_collections().collections
        if not any(c.name == COLLECTION_NAME for c in collections):
            qdrant_client.create_collection(
                collection_name=COLLECTION_NAME,
                vectors_config=VectorParams(size=vector_size, distance=Distance.COSINE)
            )
            print(f"[Worker] Created Qdrant collection: {COLLECTION_NAME} (dim={vector_size})")
    except Exception as e:
        print(f"[Worker] Error creating collection: {e}")


def download_file_from_minio(source_uri: str) -> bytes:
    if source_uri.startswith("s3://"):
        parts = source_uri[5:].split("/", 1)
        bucket = parts[0]
        object_name = parts[1] if len(parts) > 1 else ""
    else:
        raise ValueError(f"Invalid source URI: {source_uri}")
    response = minio_client.get_object(bucket, object_name)
    content = response.read()
    response.close()
    response.release_conn()
    return content


def extract_text(content: bytes, mime_type: str) -> str:
    if mime_type in ["text/plain", "text/markdown"]:
        return content.decode("utf-8", errors="ignore")

    elif mime_type == "application/pdf":
        try:
            import pdfplumber
            pages_text = []
            with pdfplumber.open(BytesIO(content)) as pdf:
                for page in pdf.pages:
                    text = page.extract_text() or ""
                    tables = page.extract_tables()
                    for table in tables:
                        for row in table:
                            cells = [str(c) if c else "" for c in row]
                            text += "\n" + " | ".join(cells)
                    pages_text.append(text)
            return "\n\n".join(pages_text)
        except Exception as e:
            print(f"[Worker] pdfplumber error, falling back to pypdf: {e}")
            try:
                from pypdf import PdfReader
                reader = PdfReader(BytesIO(content))
                return "".join(page.extract_text() or "" for page in reader.pages)
            except Exception as e2:
                print(f"[Worker] pypdf fallback error: {e2}")
                return ""

    elif mime_type in ["text/html", "application/xhtml+xml"]:
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(content, "lxml")
            for tag in soup(["script", "style", "nav", "footer", "header"]):
                tag.decompose()
            return soup.get_text(separator="\n", strip=True)
        except Exception as e:
            print(f"[Worker] HTML extraction error: {e}")
            return content.decode("utf-8", errors="ignore")

    elif mime_type in [
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword"
    ]:
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(BytesIO(content))
            parts = []
            for para in doc.paragraphs:
                if para.style and para.style.name and para.style.name.startswith("Heading"):
                    level = para.style.name.replace("Heading ", "").replace("Heading", "1")
                    try:
                        level = int(level)
                    except ValueError:
                        level = 1
                    parts.append(f"{'#' * level} {para.text}")
                elif para.text.strip():
                    parts.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    parts.append(" | ".join(cells))
            return "\n\n".join(parts)
        except Exception as e:
            print(f"[Worker] DOCX extraction error: {e}")
            return ""

    elif mime_type in [
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel"
    ]:
        try:
            from openpyxl import load_workbook
            wb = load_workbook(BytesIO(content), read_only=True, data_only=True)
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                parts.append(f"## {sheet_name}")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        parts.append(" | ".join(cells))
            wb.close()
            return "\n\n".join(parts)
        except Exception as e:
            print(f"[Worker] XLSX extraction error: {e}")
            return ""

    elif mime_type == "text/csv":
        try:
            import csv
            import io
            text = content.decode("utf-8", errors="ignore")
            reader = csv.reader(io.StringIO(text))
            rows = [" | ".join(row) for row in reader]
            return "\n".join(rows)
        except Exception as e:
            print(f"[Worker] CSV extraction error: {e}")
            return content.decode("utf-8", errors="ignore")

    elif mime_type in [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint"
    ]:
        try:
            from pptx import Presentation
            prs = Presentation(BytesIO(content))
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                slide_text = [f"## 幻灯片 {i}"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                slide_text.append(para.text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            slide_text.append(" | ".join(cells))
                parts.append("\n".join(slide_text))
            return "\n\n".join(parts)
        except Exception as e:
            print(f"[Worker] PPTX extraction error: {e}")
            return ""

    else:
        return content.decode("utf-8", errors="ignore")


def chunk_text(text: str, max_chunk_size: int = 1000, overlap: int = 100) -> List[Dict[str, Any]]:
    def _make_chunk(chunk_text_value: str, index: int) -> Dict[str, Any]:
        return {
            "index": index,
            "text": chunk_text_value,
            "text_hash": hashlib.sha256(chunk_text_value.encode()).hexdigest()[:16]
        }

    def _split_long_text(value: str) -> List[str]:
        value = value.strip()
        if not value:
            return []
        if len(value) <= max_chunk_size:
            return [value]

        pieces: List[str] = []
        start = 0
        step = max(max_chunk_size - overlap, 1)
        while start < len(value):
            end = min(start + max_chunk_size, len(value))
            pieces.append(value[start:end].strip())
            if end >= len(value):
                break
            start += step
        return [piece for piece in pieces if piece]

    chunks = []
    paragraphs = text.split("\n\n")
    current_chunk = ""
    chunk_index = 0

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue

        para_segments = _split_long_text(para)
        for segment in para_segments:
            separator = "\n\n" if current_chunk else ""
            if len(current_chunk) + len(separator) + len(segment) <= max_chunk_size:
                current_chunk += separator + segment
                continue

            if current_chunk:
                chunks.append(_make_chunk(current_chunk, chunk_index))
                chunk_index += 1
                overlap_text = current_chunk[-overlap:] if len(current_chunk) > overlap else current_chunk
                current_chunk = overlap_text + ("\n\n" if overlap_text else "") + segment
            else:
                current_chunk = segment

            if len(current_chunk) > max_chunk_size:
                overflow_segments = _split_long_text(current_chunk)
                for overflow_segment in overflow_segments[:-1]:
                    chunks.append(_make_chunk(overflow_segment, chunk_index))
                    chunk_index += 1
                current_chunk = overflow_segments[-1]

    if current_chunk:
        chunks.append(_make_chunk(current_chunk, chunk_index))
    return chunks


def get_embeddings(texts: List[str], batch_size: int = 20) -> List[List[float]]:
    model_gateway_url = os.getenv("MODEL_GATEWAY_URL", "http://model-gateway:8000")
    all_embeddings: List[List[float]] = []
    import time as _time

    for start in range(0, len(texts), batch_size):
        batch = texts[start:start + batch_size]
        max_retries = 3
        for attempt in range(max_retries):
            try:
                with httpx.Client(timeout=120.0) as client:
                    response = client.post(
                        f"{model_gateway_url}/internal/embeddings",
                        json={"texts": batch},
                    )
                    response.raise_for_status()
                    all_embeddings.extend(response.json()["embeddings"])
                    break
            except Exception as e:
                print(f"[Worker] Embedding batch [{start}:{start+len(batch)}] attempt {attempt+1} failed: {e}")
                if attempt < max_retries - 1:
                    _time.sleep(3 * (attempt + 1))
                else:
                    raise
        if start + batch_size < len(texts):
            print(f"[Worker] Embedded {start + len(batch)}/{len(texts)} chunks")

    return all_embeddings


def index_to_qdrant(document_id: str, document_title: str, source_type: str,
                    source_uri: str, chunks: List[Dict], embeddings: List[List[float]],
                    knowledge_base_id: str = None):
    if not embeddings:
        return
    ensure_collection_exists(len(embeddings[0]))
    points = []
    for chunk, embedding in zip(chunks, embeddings):
        point_id = hashlib.md5(f"{document_id}_{chunk['index']}".encode()).hexdigest()
        payload = {
            "chunk_id": point_id,
            "db_chunk_id": chunk.get("db_chunk_id", ""),
            "document_id": str(document_id),
            "title": document_title,
            "source_type": source_type,
            "source_uri": source_uri,
            "chunk_index": chunk["index"],
            "text": chunk["text"],
            "text_hash": chunk["text_hash"],
        }
        if knowledge_base_id:
            payload["knowledge_base_id"] = knowledge_base_id
        points.append(PointStruct(
            id=point_id,
            vector=embedding,
            payload=payload,
        ))
    if points:
        qdrant_client.upsert(collection_name=COLLECTION_NAME, points=points)
        print(f"[Worker] Indexed {len(points)} chunks to Qdrant (kb={knowledge_base_id})")


def delete_qdrant_vectors(document_id: str):
    """Delete all vectors for a document from Qdrant."""
    try:
        from qdrant_client.models import Filter, FieldCondition, MatchValue
        qdrant_client.delete(
            collection_name=COLLECTION_NAME,
            points_selector=Filter(
                must=[FieldCondition(key="document_id", match=MatchValue(value=document_id))]
            )
        )
        print(f"[Worker] Deleted Qdrant vectors for document: {document_id}")
    except Exception as e:
        print(f"[Worker] Error deleting Qdrant vectors: {e}")


def _run_pipeline(document, task, session):
    """Shared parse→chunk→embed→index pipeline."""
    mark_parse_processing(document)
    task.status = "running"
    task.progress = 10
    session.commit()

    # Download
    print(f"[Worker] Downloading: {document.title}")
    content = download_file_from_minio(document.source_uri)
    task.progress = 20
    session.commit()

    # Extract text
    print(f"[Worker] Extracting text")
    text = extract_text(content, document.mime_type or "text/plain")
    print(f"[Worker] Extracted {len(text)} characters")
    task.progress = 40
    session.commit()

    # Chunk
    print(f"[Worker] Chunking text")
    chunks = chunk_text(text)
    print(f"[Worker] Created {len(chunks)} chunks")
    task.progress = 50
    session.commit()

    if not chunks:
        raise ValueError("No readable content extracted from document")

    # Delete old vectors if reindexing
    delete_qdrant_vectors(str(document.id))

    # Save chunks to DB and collect UUIDs for source mapping
    db_chunks = []
    for chunk_data in chunks:
        chunk = Chunk(
            document_id=document.id,
            chunk_index=chunk_data["index"],
            text=chunk_data["text"],
            text_hash=chunk_data["text_hash"]
        )
        session.add(chunk)
        db_chunks.append(chunk)
    session.commit()

    # Enrich chunk dicts with DB-generated UUIDs (chunk↔source mapping)
    for chunk_data, db_chunk in zip(chunks, db_chunks):
        session.refresh(db_chunk)
        chunk_data["db_chunk_id"] = str(db_chunk.id)

    task.progress = 60
    session.commit()

    # Embeddings
    print(f"[Worker] Generating embeddings")
    chunk_texts = [c["text"] for c in chunks]
    embeddings = get_embeddings(chunk_texts)
    print(f"[Worker] Generated {len(embeddings)} embeddings")
    task.progress = 80
    session.commit()

    # Index to Qdrant with full source mapping
    kb_id = str(document.knowledge_base_id) if document.knowledge_base_id else None
    print(f"[Worker] Indexing to Qdrant (knowledge_base_id={kb_id})")
    index_to_qdrant(
        document_id=str(document.id),
        document_title=document.title,
        source_type=document.source_type,
        source_uri=document.source_uri,
        chunks=chunks,
        embeddings=embeddings,
        knowledge_base_id=kb_id,
    )

    task.progress = 100
    task.status = "success"
    mark_parse_completed(document)
    session.commit()

    celery_app.send_task("worker.generate_document_summary", args=[str(document.id)], queue="celery")
    celery_app.send_task("worker.extract_document_entities", args=[str(document.id)], queue="celery")

    return {"status": "success", "document_id": str(document.id), "chunks": len(chunks)}


@celery_app.task(name="worker.process_document", bind=True)
def process_document(self, document_id: str):
    """Process an uploaded document: parse → chunk → embed → index."""
    print(f"[Worker] Processing document: {document_id}")
    session = next(db.get_session())

    try:
        document = session.query(Document).filter(Document.id == document_id).first()
        if not document:
            return {"status": "error", "message": "Document not found"}

        mark_parse_processing(document)
        session.commit()

        task = session.query(IngestionTask).filter(
            IngestionTask.document_id == document_id
        ).order_by(IngestionTask.created_at.desc()).first()
        if not task:
            return {"status": "error", "message": "No ingestion task found"}

        result = _run_pipeline(document, task, session)
        print(f"[Worker] Document processed successfully: {document_id}")
        return result

    except Exception as e:
        print(f"[Worker] Error processing document {document_id}: {str(e)}")
        import traceback
        traceback.print_exc()
        if "document" in locals() and document is not None:
            mark_document_failed(document, stage="parse")
        task = session.query(IngestionTask).filter(
            IngestionTask.document_id == document_id
        ).order_by(IngestionTask.created_at.desc()).first()
        if task:
            task.status = "failed"
            task.error_message = str(e)[:500]
        session.commit()
        return {"status": "error", "message": str(e)}
    finally:
        session.close()


@celery_app.task(name="worker.fetch_and_process_url", bind=True)
def fetch_and_process_url(self, document_id: str):
    """Fetch a URL, store HTML snapshot in MinIO, then run pipeline."""
    print(f"[Worker] Fetching URL for document: {document_id}")
    session = next(db.get_session())

    try:
        document = session.query(Document).filter(Document.id == document_id).first()
        if not document:
            return {"status": "error", "message": "Document not found"}

        mark_parse_processing(document)
        session.commit()

        task = session.query(IngestionTask).filter(
            IngestionTask.document_id == document_id
        ).order_by(IngestionTask.created_at.desc()).first()
        if not task:
            return {"status": "error", "message": "No ingestion task found"}

        task.status = "running"
        task.progress = 5
        session.commit()

        # Fetch URL with safety limits
        url = document.source_uri
        max_bytes = settings.url_fetch_max_mb * 1024 * 1024
        timeout = settings.url_fetch_timeout_seconds
        max_redirects = int(os.getenv("URL_FETCH_MAX_REDIRECTS", "5"))

        print(f"[Worker] Fetching: {url}")
        with httpx.Client(
            timeout=timeout,
            follow_redirects=True,
            max_redirects=max_redirects
        ) as client:
            response = client.get(url, headers={
                "User-Agent": "CypherGuard-KB-Bot/1.0"
            })
            response.raise_for_status()

        content = response.content
        if len(content) > max_bytes:
            raise ValueError(f"Page too large: {len(content)} bytes (max {max_bytes})")

        content_type = response.headers.get("content-type", "text/html")
        if "html" in content_type:
            document.mime_type = "text/html"
        elif "text" in content_type:
            document.mime_type = "text/plain"
        else:
            document.mime_type = "text/html"

        # Store HTML snapshot to MinIO
        object_name = f"{document.owner_user_id}/url_snapshots/{document.id}.html"
        minio_client.put_object(
            settings.minio_bucket,
            object_name,
            data=BytesIO(content),
            length=len(content),
            content_type=document.mime_type
        )
        document.source_uri = f"s3://{settings.minio_bucket}/{object_name}"
        session.commit()
        print(f"[Worker] URL fetched, {len(content)} bytes stored to MinIO")

        task.task_type = "parse"
        task.progress = 15
        session.commit()

        result = _run_pipeline(document, task, session)
        print(f"[Worker] URL document processed successfully: {document_id}")
        return result

    except Exception as e:
        print(f"[Worker] Error fetching URL for {document_id}: {str(e)}")
        import traceback
        traceback.print_exc()
        if "document" in locals() and document is not None:
            mark_document_failed(document, stage="parse")
        task = session.query(IngestionTask).filter(
            IngestionTask.document_id == document_id
        ).order_by(IngestionTask.created_at.desc()).first()
        if task:
            task.status = "failed"
            task.error_message = str(e)[:500]
        session.commit()
        return {"status": "error", "message": str(e)}
    finally:
        session.close()


def generate_summary(text: str, max_chars: int = 3000) -> str:
    """Use LLM to generate document summary."""
    truncated = text[:max_chars]
    model_gateway_url = os.getenv("MODEL_GATEWAY_URL", "http://model-gateway:8000")
    try:
        with httpx.Client(timeout=120.0) as client:
            response = client.post(f"{model_gateway_url}/internal/chat", json={
                "messages": [
                    {"role": "system", "content": "你是文档摘要助手。请用中文对以下文档内容生成简洁的摘要，不超过300字。"},
                    {"role": "user", "content": truncated}
                ],
                "max_tokens": 500
            })
            response.raise_for_status()
            return response.json()["content"]
    except Exception as e:
        print(f"[Worker] Summary generation failed: {e}")
        return text[:200] + "..."


ENTITY_PATTERNS = {
    "cve": r"CVE-\d{4}-\d{4,}",
    "ipv4": r"\b(?:\d{1,3}\.){3}\d{1,3}\b",
    "domain": r"\b[a-zA-Z0-9][-a-zA-Z0-9]*(?:\.[a-zA-Z]{2,})+\b",
    "md5": r"\b[a-fA-F0-9]{32}\b",
    "sha256": r"\b[a-fA-F0-9]{64}\b",
    "email": r"\b[\w.-]+@[\w.-]+\.\w+\b",
}


def extract_entities(text: str, document_id: str, session) -> int:
    """Extract cybersecurity entities from text using regex."""
    import re
    from collections import Counter

    total = 0
    for etype, pattern in ENTITY_PATTERNS.items():
        matches = re.findall(pattern, text)
        counts = Counter(matches)
        for value, count in counts.items():
            if etype == "ipv4" and value.startswith("0.") or value.startswith("255."):
                continue
            entity = Entity(
                document_id=document_id,
                entity_type=etype,
                value=value,
                count=count
            )
            session.add(entity)
            total += 1
    if total:
        session.commit()
        print(f"[Worker] Extracted {total} entities from document {document_id}")
    return total


@celery_app.task(name="worker.generate_document_summary", bind=True)
def generate_document_summary(self, document_id: str):
    """Generate summary for a processed document."""
    session = next(db.get_session())
    try:
        document = session.query(Document).filter(Document.id == document_id).first()
        if not document:
            return {"status": "error", "message": "Document not found"}

        mark_summary_processing(document)
        session.commit()

        chunks = session.query(Chunk).filter(
            Chunk.document_id == document_id
        ).order_by(Chunk.chunk_index).limit(5).all()

        combined_text = "\n\n".join(c.text for c in chunks if c.text)
        if not combined_text:
            mark_summary_completed(document)
            session.commit()
            return {"status": "skip", "message": "No text to summarize"}

        summary = generate_summary(combined_text)
        document.summary = summary
        mark_summary_completed(document)
        session.commit()

        print(f"[Worker] Summary generated for document: {document_id}")
        return {"status": "success", "summary_length": len(summary)}
    except Exception as e:
        print(f"[Worker] Summary generation error: {e}")
        if "document" in locals() and document is not None:
            mark_document_failed(document, stage="summary")
            session.commit()
        return {"status": "error", "message": str(e)}
    finally:
        session.close()


@celery_app.task(name="worker.extract_document_entities", bind=True)
def extract_document_entities(self, document_id: str):
    """Extract cybersecurity entities from a document's chunks."""
    session = next(db.get_session())
    try:
        chunks = session.query(Chunk).filter(Chunk.document_id == document_id).all()
        combined_text = "\n".join(c.text for c in chunks if c.text)
        if not combined_text:
            return {"status": "skip"}

        session.query(Entity).filter(Entity.document_id == document_id).delete()
        session.commit()

        count = extract_entities(combined_text, document_id, session)
        return {"status": "success", "entities": count}
    except Exception as e:
        print(f"[Worker] Entity extraction error: {e}")
        return {"status": "error", "message": str(e)}
    finally:
        session.close()


@celery_app.task(name="worker.index_faq", bind=True)
def index_faq(self, faq_id: str):
    """Vectorize FAQ question and similar questions, store in Qdrant."""
    session = next(db.get_session())
    try:
        faq = session.query(FAQEntry).filter(FAQEntry.id == faq_id).first()
        if not faq:
            return {"status": "error", "message": "FAQ not found"}

        texts_to_embed = [faq.question]
        if faq.similar_questions:
            texts_to_embed.extend(faq.similar_questions)

        embeddings = get_embeddings(texts_to_embed)
        if not embeddings:
            return {"status": "error", "message": "No embeddings generated"}
        ensure_collection_exists(len(embeddings[0]))

        points = []
        for i, (text, emb) in enumerate(zip(texts_to_embed, embeddings)):
            point_id = hashlib.md5(f"faq_{faq_id}_{i}".encode()).hexdigest()
            points.append(PointStruct(
                id=point_id,
                vector=emb,
                payload={
                    "chunk_id": point_id,
                    "document_id": f"faq_{faq_id}",
                    "title": f"FAQ: {faq.question[:50]}",
                    "source_type": "faq",
                    "source_uri": "",
                    "chunk_index": i,
                    "text": f"问题: {faq.question}\n\n回答: {faq.answer}",
                    "text_hash": hashlib.sha256(text.encode()).hexdigest()[:16],
                    "faq_id": str(faq_id),
                }
            ))

        if points:
            qdrant_client.upsert(collection_name=COLLECTION_NAME, points=points)
        print(f"[Worker] Indexed FAQ: {faq_id} ({len(points)} vectors)")
        return {"status": "success", "vectors": len(points)}
    except Exception as e:
        print(f"[Worker] FAQ indexing error: {e}")
        return {"status": "error", "message": str(e)}
    finally:
        session.close()


@celery_app.task(name="worker.reindex_chunk", bind=True)
def reindex_chunk(self, chunk_id: str):
    """Re-vectorize a single chunk after edit."""
    session = next(db.get_session())
    try:
        chunk = session.query(Chunk).filter(Chunk.id == chunk_id).first()
        if not chunk or not chunk.text:
            return {"status": "error", "message": "Chunk not found or empty"}

        doc = session.query(Document).filter(Document.id == chunk.document_id).first()
        if not doc:
            return {"status": "error", "message": "Document not found"}

        embeddings = get_embeddings([chunk.text])
        if not embeddings:
            return {"status": "error", "message": "No embeddings generated"}
        ensure_collection_exists(len(embeddings[0]))

        point_id = hashlib.md5(f"{chunk.document_id}_{chunk.chunk_index}".encode()).hexdigest()
        payload = {
            "chunk_id": point_id,
            "db_chunk_id": str(chunk.id),
            "document_id": str(chunk.document_id),
            "title": doc.title,
            "source_type": doc.source_type,
            "source_uri": doc.source_uri,
            "chunk_index": chunk.chunk_index,
            "text": chunk.text,
            "text_hash": chunk.text_hash or "",
        }
        if doc.knowledge_base_id:
            payload["knowledge_base_id"] = str(doc.knowledge_base_id)
        point = PointStruct(id=point_id, vector=embeddings[0], payload=payload)
        qdrant_client.upsert(collection_name=COLLECTION_NAME, points=[point])
        print(f"[Worker] Re-indexed chunk: {chunk_id}")
        return {"status": "success"}
    except Exception as e:
        print(f"[Worker] Chunk reindex error: {e}")
        return {"status": "error", "message": str(e)}
    finally:
        session.close()


@celery_app.task(name="worker.test_task")
def test_task():
    print("[Worker] Test task running!")
    return "Task completed"


if __name__ == "__main__":
    celery_app.start()
